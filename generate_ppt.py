"""
generate_ppt.py
---------------
Generates the RedRob AI Ranker presentation as a .pptx file.
Run: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLOUR PALETTE ────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x0F, 0x1A)   # #0d0f1a  main bg
BG_CARD       = RGBColor(0x11, 0x18, 0x27)   # #111827  card bg
PURPLE        = RGBColor(0x63, 0x66, 0xF1)   # #6366f1  primary accent
PURPLE_LIGHT  = RGBColor(0xA7, 0x8B, 0xFA)   # #a78bfa  secondary accent
PURPLE_PALE   = RGBColor(0xE0, 0xE7, 0xFF)   # #e0e7ff  heading light
GREEN         = RGBColor(0x22, 0xC5, 0x5E)   # #22c55e  fresher green
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b  gold
RED_SOFT      = RGBColor(0xEF, 0x44, 0x44)   # #ef4444  warning red
TEXT_MAIN     = RGBColor(0xE2, 0xE8, 0xF0)   # #e2e8f0  body text
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8  muted
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BORDER        = RGBColor(0x1E, 0x25, 0x3B)   # subtle border

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── HELPERS ───────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg_rect(slide, x, y, w, h, fill_color, alpha=None):
    """Add a filled rectangle (background block)."""
    shape = slide.shapes.add_shape(1, x, y, w, h)   # 1 = rectangle
    shape.line.fill.background()
    shape.line.color.rgb = fill_color
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def accent_bar(slide, x, y, w, h=Pt(3).emu, color=PURPLE):
    """Thin coloured accent rule."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.line.width = 0
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=Pt(12), bold=False, color=TEXT_MAIN,
                align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb


def add_rich_textbox(slide, lines, x, y, w, h, default_size=Pt(11)):
    """
    lines = list of dicts:
      {text, size, bold, color, align, italic, space_before}
    """
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for spec in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if "space_before" in spec:
            p.space_before = spec["space_before"]
        run = p.add_run()
        run.text = spec.get("text", "")
        run.font.size   = spec.get("size", default_size)
        run.font.bold   = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = spec.get("color", TEXT_MAIN)
        run.font.name   = "Calibri"
    return txb


def slide_bg(slide):
    """Fill entire slide with dark background."""
    bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)


def section_header_band(slide, text, sub=""):
    """Top purple gradient band with slide title."""
    bg_rect(slide, 0, 0, SLIDE_W, Inches(1.25), BG_CARD)
    accent_bar(slide, 0, Inches(1.25), SLIDE_W, Pt(3).emu, PURPLE)
    # Big title
    add_textbox(slide, text,
                Inches(0.5), Inches(0.18), Inches(11), Inches(0.65),
                font_size=Pt(28), bold=True, color=PURPLE_PALE)
    if sub:
        add_textbox(slide, sub,
                    Inches(0.5), Inches(0.8), Inches(10), Inches(0.38),
                    font_size=Pt(12), color=TEXT_MUTED, italic=True)


def card(slide, x, y, w, h):
    """Rounded-looking card background."""
    shape = bg_rect(slide, x, y, w, h, BORDER)
    inner = bg_rect(slide, x+Pt(1.5).emu, y+Pt(1.5).emu,
                    w-Pt(3).emu, h-Pt(3).emu, BG_CARD)
    return inner


def bullet_block(slide, title, items, x, y, w, h,
                 title_color=PURPLE_LIGHT, bullet="▸"):
    card(slide, x, y, w, h)
    lines = [{"text": title, "size": Pt(13), "bold": True,
               "color": title_color, "align": PP_ALIGN.LEFT}]
    for item in items:
        lines.append({"text": f"{bullet}  {item}", "size": Pt(10.5),
                       "color": TEXT_MAIN, "align": PP_ALIGN.LEFT,
                       "space_before": Pt(3).emu})
    add_rich_textbox(slide, lines, x+Inches(0.15), y+Inches(0.1),
                     w-Inches(0.3), h-Inches(0.15))


def metric_box(slide, value, label, x, y, w=Inches(2.3), h=Inches(1.2),
               val_color=PURPLE_LIGHT):
    card(slide, x, y, w, h)
    add_textbox(slide, value,
                x, y+Inches(0.08), w, Inches(0.65),
                font_size=Pt(28), bold=True, color=val_color,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                x, y+Inches(0.7), w, Inches(0.4),
                font_size=Pt(9.5), color=TEXT_MUTED,
                align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """TITLE SLIDE"""
    s = blank_slide(prs)
    slide_bg(s)

    # Purple left panel
    bg_rect(s, 0, 0, Inches(5.6), SLIDE_H, BG_CARD)
    accent_bar(s, Inches(5.6), 0, Pt(4).emu, SLIDE_H, PURPLE)

    # Logo / icon area
    add_textbox(s, "🎯", Inches(0.4), Inches(0.5), Inches(1.2), Inches(1.2),
                font_size=Pt(52), align=PP_ALIGN.CENTER)

    # Main title
    add_rich_textbox(s, [
        {"text": "RedRob AI", "size": Pt(38), "bold": True,
         "color": PURPLE_PALE, "align": PP_ALIGN.LEFT},
        {"text": "Candidate Ranker", "size": Pt(38), "bold": True,
         "color": WHITE, "align": PP_ALIGN.LEFT},
    ], Inches(0.4), Inches(1.7), Inches(4.8), Inches(1.5))

    add_textbox(s, "Intelligent, explainable, fresher-friendly candidate ranking\nwithout keyword matching — and without any external APIs.",
                Inches(0.4), Inches(3.2), Inches(5.0), Inches(1.2),
                font_size=Pt(12), color=TEXT_MUTED, italic=True)

    # Right-side key stats
    stats = [
        ("89,788", "Candidates Ranked"),
        ("712/s",  "Throughput (CPU)"),
        ("30/100", "Fresher Slots"),
        ("0",      "API Calls"),
    ]
    for i, (val, lbl) in enumerate(stats):
        col = i % 2
        row = i // 2
        metric_box(s, val, lbl,
                   Inches(6.1) + col*Inches(3.5),
                   Inches(1.2) + row*Inches(1.5),
                   w=Inches(3.2), h=Inches(1.25),
                   val_color=AMBER if "0" in val else PURPLE_LIGHT)

    # Bottom tagline
    add_textbox(s, "RedRob AI Challenge  |  Dabbara Santhosh Kumar  |  github.com/Santhu7718/redrob-ai-ranker",
                Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
                font_size=Pt(8.5), color=TEXT_MUTED)


def slide_02_solution_overview(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Solution Overview",
                        "What was built · How it differs from traditional systems")

    # Left col — what it is
    bullet_block(s,
        "What We Built",
        [
            "A rule-based multi-signal hybrid scoring engine",
            "Ranks 89,788 candidates in ~126 seconds on CPU alone",
            "Fully offline — zero API calls, zero GPU required",
            "Streamlit web UI: upload any file → get ranked shortlist",
            "Dual-track system guarantees fresher representation",
            "Every score is 100% auditable and explainable",
        ],
        Inches(0.35), Inches(1.45), Inches(6.1), Inches(3.2),
        title_color=PURPLE_LIGHT)

    # Right col — differentiators
    bullet_block(s,
        "What Makes This Different",
        [
            "Not keyword matching — trust-weighted skill scoring",
            "Endorsement + duration required to claim a skill",
            "Career text-mining finds implicit expertise",
            "YoE curve penalises over-experience, rewards potential",
            "Services firm penalty (TCS/Infosys with no real ML work)",
            "Behavioral signals: inactive candidates deprioritised",
        ],
        Inches(6.7), Inches(1.45), Inches(6.25), Inches(3.2),
        title_color=GREEN)

    # Bottom comparison row
    bg_rect(s, Inches(0.35), Inches(4.85), Inches(12.6), Inches(0.04), BORDER)
    labels = [
        ("Traditional ATS", "Keywords only · Binary match · No reasoning", RED_SOFT),
        ("Vector Search",   "Semantic match · Still no signal fusion · Black box", AMBER),
        ("Our Approach",    "6 signals · Trust-weighted · Fresher uplift · Explainable", GREEN),
    ]
    for i, (title, desc, col) in enumerate(labels):
        x = Inches(0.35) + i * Inches(4.3)
        card(s, x, Inches(4.95), Inches(4.1), Inches(1.35))
        accent_bar(s, x, Inches(4.95), Inches(4.1), Pt(3).emu, col)
        add_textbox(s, title, x+Inches(0.12), Inches(5.06),
                    Inches(3.9), Inches(0.35),
                    font_size=Pt(11), bold=True, color=col)
        add_textbox(s, desc, x+Inches(0.12), Inches(5.4),
                    Inches(3.9), Inches(0.7),
                    font_size=Pt(9.5), color=TEXT_MUTED)


def slide_03_jd_and_signals(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "JD Understanding & Candidate Evaluation",
                        "What the role demands · Which signals matter most")

    # JD Skills block
    bullet_block(s,
        "Key Requirements Extracted from JD",
        [
            "55 critical skills parsed (Python · PyTorch · TensorFlow · NLP · BERT · Transformers)",
            "32 preferred skills (FastAPI · Docker · Kubernetes · Spark · SQL)",
            "High-value ML skills: LoRA · QLoRA · PEFT · RAG · FAISS · Qdrant · Pinecone",
            "Location signal: India-based preferred, no hard block",
            "Implicit red flags: keyword-stuffed CVs, TCS-style zero-ML experience",
            "Extracted via regex + curated ontology from job_description.docx — no LLM",
        ],
        Inches(0.35), Inches(1.45), Inches(6.1), Inches(2.9),
        title_color=AMBER)

    # Candidate signals block
    bullet_block(s,
        "Candidate Signals Beyond Keywords",
        [
            "Trust-weighted skill score: proficiency × endorsements × duration",
            "Career trajectory: promotion speed, company tier, AI-specific roles",
            "Platform signals: GitHub activity, profile completeness, response rate",
            "Education: institution tier (IIT/NIT/BITS) + field + GPA",
            "Certifications: ML-relevant courses (DeepLearning.AI, Kaggle, NVIDIA)",
            "Behavioral: last-active days, open-to-work flag, availability",
        ],
        Inches(6.7), Inches(1.45), Inches(6.25), Inches(2.9),
        title_color=PURPLE_LIGHT)

    # Trust formula
    card(s, Inches(0.35), Inches(4.55), Inches(12.6), Inches(1.6))
    add_textbox(s, "🔑  Anti-Keyword-Stuffing Formula",
                Inches(0.55), Inches(4.65), Inches(6), Inches(0.35),
                font_size=Pt(12), bold=True, color=AMBER)
    add_textbox(s,
        'skill_trust  =  0.50 × proficiency_score  +  0.30 × min(1, endorsements÷20)  +  0.20 × min(1, months_used÷12)',
                Inches(0.55), Inches(5.05), Inches(8.5), Inches(0.35),
                font_size=Pt(10.5), color=PURPLE_PALE, bold=True)
    add_textbox(s,
        "A candidate claiming 'FAISS' with 0 endorsements and 0 months gets near-zero credit — regardless of how many times it appears in their profile.",
                Inches(0.55), Inches(5.45), Inches(11), Inches(0.4),
                font_size=Pt(9.5), color=TEXT_MUTED, italic=True)


def slide_04_ranking_methodology(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Ranking Methodology",
                        "How candidates are retrieved, scored, and ranked")

    # 6 signals
    signals = [
        ("🎯  Skill Match",      "35%", "Trust-weighted coverage of JD skills + high-value ML boost", PURPLE_LIGHT),
        ("📅  Experience Fit",   "25%", "YoE curve: 0→0.65  2→0.72  5→0.85  12+→0.60", GREEN),
        ("🎓  Education",        "15%", "Tier (IIT+0.25) · Field (CS/ML+0.15) · GPA · Degree level", AMBER),
        ("📋  Completeness",     "10%", "Fields filled · GitHub / LinkedIn / portfolio presence", PURPLE_LIGHT),
        ("📜  Certifications",   "10%", "ML-relevant upskilling: Coursera · Kaggle · NVIDIA · HuggingFace", GREEN),
        ("🔍  Keyword Density",   "5%", "Full-text implicit JD mentions across entire profile", AMBER),
    ]
    for i, (name, weight, desc, col) in enumerate(signals):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.35) + col_idx * Inches(6.45)
        y = Inches(1.52) + row_idx * Inches(1.28)
        card(s, x, y, Inches(6.2), Inches(1.18))
        accent_bar(s, x, y, Inches(6.2), Pt(3).emu, col)

        add_textbox(s, name, x+Inches(0.15), y+Inches(0.1),
                    Inches(3.8), Inches(0.35), font_size=Pt(11.5),
                    bold=True, color=col)
        add_textbox(s, weight, x+Inches(4.5), y+Inches(0.1),
                    Inches(1.5), Inches(0.35), font_size=Pt(14),
                    bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        add_textbox(s, desc, x+Inches(0.15), y+Inches(0.5),
                    Inches(5.9), Inches(0.55), font_size=Pt(9.5),
                    color=TEXT_MUTED)

    # Formula + dual-track
    card(s, Inches(0.35), Inches(5.4), Inches(8.2), Inches(1.75))
    add_textbox(s, "Composite Formula",
                Inches(0.5), Inches(5.5), Inches(4), Inches(0.3),
                font_size=Pt(10), bold=True, color=PURPLE_LIGHT)
    add_textbox(s,
        "score = Σ(signal_i × weight_i)  ×  fresher_uplift (≤ 1.30×)\n\nAll signals normalised 0→1. Dual-track: top 70 experienced + top 30 freshers merged & re-sorted.",
                Inches(0.5), Inches(5.85), Inches(7.9), Inches(1.1),
                font_size=Pt(9.5), color=TEXT_MAIN)

    card(s, Inches(8.75), Inches(5.4), Inches(4.2), Inches(1.75))
    add_textbox(s, "No Models / No Training",
                Inches(8.9), Inches(5.5), Inches(4), Inches(0.3),
                font_size=Pt(10), bold=True, color=GREEN)
    add_textbox(s,
        "Pure Python + NumPy + regex.\nDeterministic, auditable, reproducible.\nSame input → same output, always.",
                Inches(8.9), Inches(5.85), Inches(3.9), Inches(1.1),
                font_size=Pt(9.5), color=TEXT_MUTED)


def slide_05_explainability(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Explainability & Data Validation",
                        "How decisions are explained · How bad data is handled")

    bullet_block(s,
        "Per-Candidate Reasoning",
        [
            "Every candidate gets a structured reasoning string output to CSV",
            "Format: SKILLS: matched 12/55 JD skills: pytorch, lora... || EXPERIENCE: Senior (7.2 YoE) || ...",
            "6 individual signal scores surfaced alongside the composite",
            "Fresher uplift multiplier shown explicitly (e.g. ×1.28 BOOST)",
            "Streamlit UI shows bar-chart breakdown per signal for every rank",
            "All score components are deterministic — no randomness",
        ],
        Inches(0.35), Inches(1.45), Inches(6.1), Inches(3.0),
        title_color=PURPLE_LIGHT)

    bullet_block(s,
        "Anti-Hallucination & Anti-Gaming",
        [
            "No LLM involved → zero hallucination risk in reasoning",
            "Skills require endorsements + duration to count (trust multiplier)",
            "Keyword-stuffed profiles are explicitly penalised",
            "Services-firm penalty: TCS/Infosys/Wipro bulk contractors → ×0.60",
            "Behavioral availability check: inactive 180d + low response → ×0.50",
            "Suspicious profiles (short stint spam, mismatched title) get penalties",
        ],
        Inches(6.7), Inches(1.45), Inches(6.25), Inches(3.0),
        title_color=RED_SOFT)

    # Anti-gaming table
    penalties = [
        ("Services firm, no ML content",     "×0.60", RED_SOFT),
        ("Inactive >180 days + low response", "×0.50", RED_SOFT),
        ("Skill claimed but 0 endorsements",  "×0.30 trust", AMBER),
        ("Fresher with strong signals",       "×1.10–1.30", GREEN),
        ("High-value ML skill match bonus",   "+0.05 per hit", GREEN),
    ]
    card(s, Inches(0.35), Inches(4.65), Inches(12.6), Inches(2.55))
    add_textbox(s, "Penalty & Boost Reference Table",
                Inches(0.55), Inches(4.75), Inches(6), Inches(0.3),
                font_size=Pt(11), bold=True, color=AMBER)
    for i, (condition, mult, col) in enumerate(penalties):
        y_off = Inches(5.15) + i * Inches(0.37)
        add_textbox(s, f"• {condition}", Inches(0.55), y_off, Inches(9), Inches(0.32),
                    font_size=Pt(9.5), color=TEXT_MAIN)
        add_textbox(s, mult, Inches(11.5), y_off, Inches(1.3), Inches(0.32),
                    font_size=Pt(9.5), bold=True, color=col, align=PP_ALIGN.RIGHT)


def slide_06_workflow(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "End-to-End Workflow",
                        "From JD input to ranked candidate output")

    steps = [
        ("1", "JD Input",         "job_description.docx parsed\nvia jd_parser.py\n55 critical + 32 preferred\nskills extracted",       PURPLE),
        ("2", "Load Candidates",  "candidates.jsonl streamed\nline-by-line\n89,788 profiles loaded\nMalformed lines skipped",           PURPLE_LIGHT),
        ("3", "Score Each",       "6-signal scorer runs\non every candidate\n712 candidates/sec\n~126s total on CPU",                   GREEN),
        ("4", "Dual-Track Split", "Freshers ≤2 YoE separated\nExperienced 70 slots\nFreshers 30 slots\nFresher uplift applied",        AMBER),
        ("5", "Merge & Rank",     "Both tracks sorted\nby final_score\nMerged → top 100\nFinal re-sort by score",                      PURPLE_LIGHT),
        ("6", "Output",           "submission.csv written\nRank · Score · Reasoning\nValidation: 100 rows ✓\nUnique ranks ✓",          GREEN),
    ]

    for i, (num, title, body, col) in enumerate(steps):
        x = Inches(0.3) + i * Inches(2.16)
        y = Inches(1.5)
        w = Inches(2.0)
        h = Inches(4.5)
        card(s, x, y, w, h)
        accent_bar(s, x, y, w, Pt(4).emu, col)

        # Circle number
        add_textbox(s, num,
                    x + Inches(0.7), y + Inches(0.12),
                    Inches(0.55), Inches(0.45),
                    font_size=Pt(20), bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title,
                    x + Inches(0.08), y + Inches(0.65),
                    Inches(1.85), Inches(0.45),
                    font_size=Pt(10.5), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, body,
                    x + Inches(0.1), y + Inches(1.18),
                    Inches(1.85), Inches(3.0),
                    font_size=Pt(9), color=TEXT_MUTED,
                    align=PP_ALIGN.CENTER)

        # Arrow between steps (not after last)
        if i < len(steps) - 1:
            add_textbox(s, "→",
                        x + Inches(2.01), y + Inches(2.1),
                        Inches(0.22), Inches(0.35),
                        font_size=Pt(16), color=PURPLE_LIGHT,
                        align=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(s,
        "UI Path: Upload file (any format) → paste JD → click Rank → view breakdown → download CSV/JSON",
                Inches(0.35), Inches(6.2), Inches(12.6), Inches(0.35),
                font_size=Pt(9.5), color=TEXT_MUTED, italic=True,
                align=PP_ALIGN.CENTER)


def slide_07_architecture(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "System Architecture",
                        "Module breakdown and data flow")

    # Three-layer architecture visual
    layers = [
        ("INPUT LAYER",
         ["job_description.docx", "candidates.jsonl (89K)", "Google Forms CSV / Excel / JSON"],
         Inches(0.3), PURPLE),
        ("PROCESSING LAYER",
         ["jd_parser.py  →  55 critical skills", "universal_parser.py  →  column auto-detect",
          "scorer.py  →  6 signals × weights", "rank.py  →  dual-track split + merge"],
         Inches(4.65), PURPLE_LIGHT),
        ("OUTPUT LAYER",
         ["submission.csv  (rank · score · reasoning)", "ranked_candidates.json",
          "Streamlit UI  →  visual breakdown", "GitHub Actions CI  →  auto-validated"],
         Inches(9.0), GREEN),
    ]

    for (label, items, x, col) in layers:
        card(s, x, Inches(1.45), Inches(4.1), Inches(4.0))
        accent_bar(s, x, Inches(1.45), Inches(4.1), Pt(4).emu, col)
        add_textbox(s, label, x+Inches(0.15), Inches(1.55),
                    Inches(3.8), Inches(0.4),
                    font_size=Pt(11), bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_textbox(s, f"▸  {item}",
                        x+Inches(0.2), Inches(2.1) + j*Inches(0.65),
                        Inches(3.7), Inches(0.55),
                        font_size=Pt(9.5), color=TEXT_MAIN)

        if x != Inches(9.0):
            add_textbox(s, "⟶",
                        x+Inches(4.12), Inches(3.2),
                        Inches(0.5), Inches(0.4),
                        font_size=Pt(20), color=PURPLE, align=PP_ALIGN.CENTER)

    # File map at bottom
    card(s, Inches(0.3), Inches(5.65), Inches(12.65), Inches(1.6))
    add_textbox(s, "Key Files",
                Inches(0.5), Inches(5.75), Inches(3), Inches(0.3),
                font_size=Pt(10), bold=True, color=PURPLE_LIGHT)
    files = [
        ("app.py",               "Streamlit UI"),
        ("universal_parser.py",  "Any-format parser"),
        ("universal_scorer.py",  "Flat-data scorer"),
        ("scorer.py",            "Challenge scorer"),
        ("jd_parser.py",         "JD skill extractor"),
        ("rank.py",              "CLI ranker + dual-track"),
        ("Dockerfile",           "Container build"),
        ("ci.yml",               "GitHub Actions CI"),
    ]
    for i, (fname, desc) in enumerate(files):
        col_i = i % 4
        row_i = i // 4
        x = Inches(0.5) + col_i * Inches(3.1)
        y = Inches(6.1) + row_i * Inches(0.45)
        add_textbox(s, f"📄 {fname}", x, y, Inches(1.8), Inches(0.35),
                    font_size=Pt(9), bold=True, color=AMBER)
        add_textbox(s, desc, x+Inches(1.85), y, Inches(1.2), Inches(0.35),
                    font_size=Pt(9), color=TEXT_MUTED)


def slide_08_results(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Results & Performance",
                        "Challenge outcomes and runtime validation")

    # Top metrics
    metrics = [
        ("89,788",  "Candidates Scored",      PURPLE_LIGHT),
        ("712/s",   "Scoring Speed (CPU)",     GREEN),
        ("~126s",   "Total Runtime",           AMBER),
        ("30/100",  "Freshers in Top 100",     GREEN),
        ("0.9068",  "Top Composite Score",     PURPLE_LIGHT),
        ("0",       "API / GPU Calls",         AMBER),
    ]
    for i, (val, lbl, col) in enumerate(metrics):
        col_i = i % 3
        row_i = i // 3
        metric_box(s, val, lbl,
                   Inches(0.35) + col_i * Inches(4.3),
                   Inches(1.45) + row_i * Inches(1.4),
                   w=Inches(4.0), h=Inches(1.2), val_color=col)

    # Top 5 candidates
    card(s, Inches(0.35), Inches(4.25), Inches(12.6), Inches(2.98))
    add_textbox(s, "🏆  Top 10 Candidates — Snapshot",
                Inches(0.55), Inches(4.35), Inches(6), Inches(0.35),
                font_size=Pt(11), bold=True, color=AMBER)

    top5 = [
        ("#1", "CAND_0081846", "Lead AI Engineer @ Razorpay",         "6.7 YoE", "0.9068"),
        ("#2", "CAND_0046064", "Senior NLP Engineer @ Salesforce",    "8.9 YoE", "0.8994"),
        ("#3", "CAND_0027691", "NLP Engineer @ Haptik",               "6.5 YoE", "0.8953"),
        ("#4", "CAND_0077337", "Staff ML Engineer @ Paytm",           "7.0 YoE", "0.8919"),
        ("#5", "CAND_0079387", "AI Engineer @ Microsoft",             "6.9 YoE", "0.8858"),
    ]
    headers = ["Rank", "Candidate ID", "Role & Company", "Experience", "Score"]
    col_xs  = [Inches(0.55), Inches(1.3), Inches(3.0), Inches(9.5), Inches(11.4)]
    col_ws  = [Inches(0.7),  Inches(1.6), Inches(6.0), Inches(1.8),  Inches(1.5)]

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        add_textbox(s, hdr, cx, Inches(4.75), cw, Inches(0.28),
                    font_size=Pt(8.5), bold=True, color=TEXT_MUTED,
                    align=PP_ALIGN.LEFT)

    for i, row_data in enumerate(top5):
        y_r = Inches(5.1) + i * Inches(0.38)
        bg_rect(s, Inches(0.45), y_r - Pt(2).emu,
                Inches(12.4), Inches(0.36),
                BG_CARD if i % 2 == 0 else BORDER)
        for j, (cell, cx, cw) in enumerate(zip(row_data, col_xs, col_ws)):
            color = AMBER if j == 0 else (GREEN if j == 4 else TEXT_MAIN)
            add_textbox(s, cell, cx, y_r, cw, Inches(0.34),
                        font_size=Pt(9), color=color, bold=(j in (0,4)))

    add_textbox(s, "Validation: python dataset/validate_submission.py submission.csv  →  ✅ Submission is valid.",
                Inches(0.55), Inches(7.05), Inches(10), Inches(0.3),
                font_size=Pt(9), color=GREEN, italic=True)


def slide_09_technologies(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Technologies Used",
                        "What was used · Why each was chosen")

    tech = [
        ("Python 3.11",       "Core language",          "Speed, ecosystem maturity, rich standard library",                PURPLE_LIGHT),
        ("NumPy",             "Numerical ops",          "Fast array math for score aggregation across 89K candidates",     PURPLE_LIGHT),
        ("pandas",            "Data handling",          "Flexible CSV/Excel ingestion with type coercion",                 PURPLE_LIGHT),
        ("regex (re)",        "Skill extraction",       "Compiled patterns for fast, accurate JD and profile scanning",    PURPLE_LIGHT),
        ("python-docx",       "JD parsing",             "Reads Word doc JD natively — no conversion needed",              AMBER),
        ("Streamlit",         "Web UI",                 "Zero-config Python-native web framework; rapid iteration",        AMBER),
        ("openpyxl / xlrd",   "Excel support",          "Read .xlsx/.xls Google Forms exports directly",                  AMBER),
        ("Docker",            "Containerisation",       "Reproducible deployment; one-command run anywhere",              GREEN),
        ("GitHub Actions",    "CI/CD",                  "Automated syntax + integration test on every push to main",      GREEN),
        ("scikit-learn",      "Utilities (optional)",   "Available for future feature; percentile normalisation",         GREEN),
    ]

    for i, (name, category, reason, col) in enumerate(tech):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.35) + col_i * Inches(6.45)
        y = Inches(1.52) + row_i * Inches(1.04)
        card(s, x, y, Inches(6.2), Inches(0.94))
        accent_bar(s, x, y, Pt(4).emu, Inches(0.94), col)
        add_textbox(s, name, x+Inches(0.2), y+Inches(0.05),
                    Inches(2.2), Inches(0.35),
                    font_size=Pt(11), bold=True, color=col)
        add_textbox(s, category, x+Inches(2.5), y+Inches(0.08),
                    Inches(3.5), Inches(0.28),
                    font_size=Pt(9), color=TEXT_MUTED, italic=True)
        add_textbox(s, reason, x+Inches(0.2), y+Inches(0.5),
                    Inches(5.9), Inches(0.38),
                    font_size=Pt(9.5), color=TEXT_MAIN)

    # NOT used section
    card(s, Inches(0.35), Inches(6.8), Inches(12.6), Inches(0.55))
    add_textbox(s, "❌  NOT used:  OpenAI · Gemini · Hugging Face Inference API · FAISS at runtime · any GPU · any external network call during ranking",
                Inches(0.55), Inches(6.88), Inches(12), Inches(0.38),
                font_size=Pt(9.5), color=RED_SOFT, italic=True)


def slide_10_submission(prs):
    s = blank_slide(prs)
    slide_bg(s)
    section_header_band(s, "Submission Assets",
                        "What is delivered and where to find it")

    assets = [
        ("📄 submission.csv",
         "Primary deliverable",
         "100 ranked candidates · rank · score · full reasoning string\nValidated: ✅ Submission is valid.",
         AMBER),
        ("🖥  app.py + UI",
         "Interactive web app",
         "Upload any file → JD → Rank → Download\nStreamlit dark UI with signal breakdowns per candidate",
         PURPLE_LIGHT),
        ("⚙️  scorer.py + rank.py",
         "CLI pipeline",
         "Full scoring engine for the 89K challenge dataset\nRuns in ~126s on CPU · 712 candidates/second",
         GREEN),
        ("🌐  universal_parser.py",
         "Any-format ingestor",
         "Reads CSV / Excel / JSON / JSONL · Auto-detects 50+ column aliases\nGoogle Forms compatible out of the box",
         PURPLE_LIGHT),
        ("🐳  Dockerfile + docker-compose",
         "Deployment",
         "docker compose up --build → live on :8501\nOne-command deploy anywhere",
         AMBER),
        ("🔁  .github/workflows/ci.yml",
         "CI / CD",
         "GitHub Actions · Syntax check + integration test\nPasses in <30 seconds on every push",
         GREEN),
    ]

    for i, (name, cat, desc, col) in enumerate(assets):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.35) + col_i * Inches(6.45)
        y = Inches(1.52) + row_i * Inches(1.7)
        card(s, x, y, Inches(6.2), Inches(1.6))
        accent_bar(s, x, y, Inches(6.2), Pt(3).emu, col)
        add_textbox(s, name, x+Inches(0.15), y+Inches(0.1),
                    Inches(4.5), Inches(0.38),
                    font_size=Pt(12), bold=True, color=col)
        add_textbox(s, cat, x+Inches(0.15), y+Inches(0.5),
                    Inches(5.9), Inches(0.25),
                    font_size=Pt(9.5), bold=True, color=TEXT_MUTED)
        add_textbox(s, desc, x+Inches(0.15), y+Inches(0.78),
                    Inches(5.9), Inches(0.7),
                    font_size=Pt(9.5), color=TEXT_MAIN)

    # GitHub link
    card(s, Inches(0.35), Inches(6.75), Inches(12.6), Inches(0.5))
    add_textbox(s,
        "🔗  GitHub:  github.com/Santhu7718/redrob-ai-ranker   |   All source code, CI badge, Docker, README, and submission.csv",
                Inches(0.55), Inches(6.83), Inches(12), Inches(0.35),
                font_size=Pt(10), bold=True, color=PURPLE_LIGHT,
                align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs)
    print("  ✓ Slide 1 — Title")
    slide_02_solution_overview(prs)
    print("  ✓ Slide 2 — Solution Overview")
    slide_03_jd_and_signals(prs)
    print("  ✓ Slide 3 — JD Understanding & Candidate Evaluation")
    slide_04_ranking_methodology(prs)
    print("  ✓ Slide 4 — Ranking Methodology")
    slide_05_explainability(prs)
    print("  ✓ Slide 5 — Explainability & Data Validation")
    slide_06_workflow(prs)
    print("  ✓ Slide 6 — End-to-End Workflow")
    slide_07_architecture(prs)
    print("  ✓ Slide 7 — System Architecture")
    slide_08_results(prs)
    print("  ✓ Slide 8 — Results & Performance")
    slide_09_technologies(prs)
    print("  ✓ Slide 9 — Technologies Used")
    slide_10_submission(prs)
    print("  ✓ Slide 10 — Submission Assets")

    out = "RedRob_AI_Ranker_Presentation.pptx"
    prs.save(out)
    print(f"\n✅  Saved → {out}")


if __name__ == "__main__":
    main()
